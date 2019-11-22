#!usr/bin/env python
import os, sys, argparse, re

from time import sleep
from abc import ABC, abstractmethod
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR


# \\\\\\\\\\\\\\\\\\\\\\\\\\\\\\
# *** GLOBAL CONSTANTS ****
# //////////////////////////////

# Path to textfile listing template paths
TXT_PATH = "paths_to_templates.txt"

# For cover slide
COVER_SLIDE = 0
TITLE_SHAPE = 0 # Same for other slides
DATE_SHAPE = 1 
CREATORS_TABLE_CONF = 1
CREATORS_TABLE_REP = 2
TABLE_COORDS = {
    "preparers": (0,1),
    "reviewers": (1,1),
    # "approvers": (2,1) 
}

# slide index of table of contents
TOC = 2


# \\\\\\\\\\\\\\\\\\\\\\\\\\\\\\
# *** FUNCTION DEFINITIONS ****
# //////////////////////////////

def fetch_interfaces(dir_path):
    """
    Crawls simulation directory to fetch names of interfaces
    """
    if not os.path.isabs(dir_path) or \
        os.path.split(dir_path)[1] != "Simulation":
        print("ERROR: Simulation folder could not be found.")
        sys.exit(-1)

    return os.listdir(dir_path)


def weave_reports(rep_type, conf_path, interfaces=[]):
    """
    Generate reports based on input confirmation tools and indicated type
    """
    ct = ConfirmationTools(conf_path) # Initialize confirmation tools
    reports = init_reports(rep_type, ct, interfaces) # Initialize reports
    for rep in reports:
        make_cover(ct, rep)
        copy_slides(ct, rep)
        save_report(rep)


def init_reports(rep_type, conf_tools, interfaces=[]):
    """
    Initializes and returns Report based on user input and template
    """
    templates = _load_template_paths(TXT_PATH)
    reports = None
    # Instantiate report based on user input
    if rep_type == "si" and interfaces:
        reports = [ SIReport(templates[rep_type], interface) for interface in interfaces ]
    elif rep_type == "pi":
        reports = PIReport(templates[rep_type])
    elif rep_type == "emc":
        reports = EMCReport(templates[rep_type])
    else:
        reports = ThermalReport(templates[rep_type])

    if not isinstance(reports, list):
        reports = list(reports)

    return reports


def _load_template_paths(file_path):
    """
    Fetches template paths and returns dict mapping report type to template
    """
    # dict to be populated
    templates = {
        "si": "",
        "pi": "",
        "emc": "",
        "thermal": "",
    }

    # Fetch paths from text file in same directory
    with open(file_path, "r") as f:
        for rep_type in templates:
            # Iterate over each line prefixed with template key and = (e.g. "si=")
            for line in f.readlines():
                if line.startswith(rep_type):
                    start = line.index("=")
                    templates[rep_type] = line[start+1:] # Omit delimiter and copy remaining str into dict
    
    return templates


def make_cover(conf_tools, report):
    """
    Generates first slide of report
    """
    cover = report.pptx.slides[COVER_SLIDE]
    title = __get_title(report) 

    cover.shapes[TITLE_SHAPE].text = title
    cover.shapes[DATE_SHAPE].text = __get_date()
    creators = conf_tools.get_creators()

    # Match table coordinates with Report.creator keys and insert values of latter
    for key, coords in TABLE_COORDS.items():
        for party in creators: 
            if key == party:
                cover.shapes[CREATORS_TABLE_REP].cell(coords[0], coords[1]).text = creators[party]

    print(f"Cover slide generated for {title}.")


def __get_title(report):
    """
    Returns user input title for a particular file
    """
    title = ""
    while True:
       title = input(f"Input title for {report}: ")
       if title:
           break

    return title


# def __is_rep_type():
#     """
#     Returns a valid report type from user
#     """
#     sim_types = ["si", "pi", "emc", "thermal"] # Types of PCB simulation reports
#     while True:
#         rep_type = input("Input type of report (SI / PI / EMC / Thermal): " )
#         # Verify user input
#         if rep_type.lower() in sim_types:
#             break

#     return rep_type


def __get_date():
    """
    Returns a user inputted date formatted into Japanese
    """
    date = ""
    while True:
        # Instructions for user input
        prompt = """
        Input report date as follows:

        yyyy,MM,dd

        where:
        yyyy -> year
        MM -> month
        dd -> date
        """
        # Check if instructions were followed
        try:
            date = input(prompt).split(",") # Split for output str formatting
            break
        except: 
            continue

    return u"{0}年　{1}月　{2}日".format(date[0], date[1], date[2])


def copy_slides(conf_tools, report):
    """
    Copies over slides of interest 
    from confirmation tools to report
    """
    toc = conf_tools.get_toc()
    # Read each slide according to toc
    for section in toc:
        for slide_nums in toc[section]:
            for slide_num in slide_nums:
                # Get curr slide of conf_tools
                conf_slide = conf_tools.pptx.slides[slide_num]
                
                # Add a slide of the same format in the report
                layout_name = conf_tools.pptx.slides[slide_num].slide_layout.name
                layout = report.pptx.slides.slide_layout.get_by_name(layout_name)
                report.pptx.slides.add_slide(layout)

                # Fetch newly appended slide
                new_index = len(report.pptx.slides) - 1
                rep_slide = report.pptx.slides[new_index]

                # Copy title of conf_tools slide to report slide
                rep_slide.shapes[TITLE_SHAPE].text = conf_slide.shapes[TITLE_SHAPE].text[:]

                __copy_shapes(conf_slide, conf_tools.slides.index(conf_slide), rep_slide)

    return report


def __copy_shapes(src_slide, src_index, dest_slide):
    """
    Iterates over source slide's shapes 
    and creates matching shapes in destination slide
    Returns None (mutates dest_slide)
    """
    curr = 1 # Shape pointer; starts at 1 to exclude title shape
    for shape in src_slide.shapes[TITLE_SHAPE+1:]:
        # Read position of original
        top = shape.top
        left = shape.left
        height = shape.height
        width = shape.width                    

        if shape.has_text_frame:
            # Add new textbox and add text thereto
            dest_slide.shapes.add_textbox(left, top, width, height)
            dest_slide.shapes[curr].text = shape.text[:]
            curr += 1

        elif shape.has_table:
            # Extract data from table
            table = shape.table
            cols = len(table.columns)
            rows = len(table.rows)
            cells = table.iter_cells() # Cell generator w/ text prop

            # Make table
            dest_slide.shapes.add_table(rows, cols, left, top, width, height)

            # Copy over contents from original to new table
            for cell in cells:
                for new_cell in dest_slide.shapes[curr].table.iter_cells():
                    new_cell.text = cell.text[:]
            curr += 1

        else:
            # Likely a group shape
            try:
                dest_slide.shapes.add_group_shape()
                __copy_group_shape(shape, dest_slide.shapes[curr])
                curr += 1

            except AttributeError: # In case not group shape
                print(f"Unidentifiable shape found in slide {src_index}")


def __copy_group_shape(src_shape, dest_shape):
    """
    Iterates over source group shape
    and creates matching shapes in destination group shape
    Returns None (mutates dest_shape)
    """
    curr_sub = 0 # To track subshapes added to group shape
    img_cnt = 1 # To differentiate img filenames
    for shape in src_shape.shapes:
        # Get position and dimensions of subshape
        sub_left = shape.left
        sub_top = shape.top
        sub_width = shape.width
        sub_height = shape.height

        if shape.shape_type == 13: # PICTURE
            # Make temp folder to extract and save imgs as files
            blob = shape.image.blob
            if not "temp" in os.listdir("."):
                os.mkdir("temp")
            filename = f"temp_img-{img_cnt}.png"
            img_path = os.path.join(os.getcwd(), "temp", filename)
            with open(img_path, "wb") as f:
                f.write(blob)
            dest_shape.shapes.add_picture(filename, sub_left, sub_top,\
                                          sub_width, sub_height)
            img_cnt += 1

        elif shape.shape_type == 17: # TEXT_BOX
            dest_shape.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
            dest_shape.shapes[curr_sub].text = shape.text
        
        elif shape.shape_type == None: 
            # Assuming shape is Connector
            begin = shape.begin_x, shape.begin_y
            end = shape.end_x, shape.end_y
            # Default to Straight
            dest_shape.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin, end)
            dest_shape
        # Move to next shape
        curr_sub += 1

    # Cleanup temp folder made for images
    if "temp" in os.listdir("."):
        os.rmdir("temp")


def save_report(report):
    """
    Gets filename from user and closes report after saving
    """
    filename = ""
    path = ""
    while True:
        filename = input("Input filename to save the report: ")
        path = input("Input path to save report: ")
        if filename and os.path.isabs(path):
            break

    report.pptx.save(os.path.join(path, filename))
    print(f"{filename} saved in {path}.")


# def make_TOU(report):
#     """Generates Table of Updates"""
#     pass


# def make_TOC(report):
#     """Generates Table of Contents"""
#     pass


# \\\\\\\\\\\\\\\\\\\\\\\\\\\
# *** CLASS DEFINITIONS ****
# ///////////////////////////

# ===============
# Report classes
# ===============

class Report(ABC):
    """
    Base class for simulation report
    """
    def __init__(self, template):
        self.__pptx = Presentation(template)
    
    @property
    def pptx(self):
        return self.__pptx


class ConfirmationTools(Report):
    """
    Class for initial, pre-simulation report
    """
    def __init__(self, file):
        super().__init__(file)

    def get_creators(self):
        """
        Gets list of authors, reviewers, and approvers 
        from Confirmation Tools object
        """
        creators = {
            "preparers": "",
            "reviewers": "",
            # "approvers": ""
        }

        for party, coords in TABLE_COORDS.items():
            if party in creators.keys():
                creators[party] = self.pptx.slides[COVER_SLIDE].\
                                shapes[CREATORS_TABLE_CONF].table.cell(coords[0], coords[1]).text[:]

        return creators

    def get_toc(self):
        """
        Returns dict of section->slide_num(s) for sections of interest
        """
        toc = self.pptx.slides[TOC].shapes[0].table # Save table
        # To be populated with slide nums
        toc_dict = {
            "sim_targets": None,
            "eye_masks": None,
            "topology": None
        }

        row = 1 # Starting y coord of table traversal
        while True:
            section_name = toc.cell(row, 0).text[:]
            section_name = section_name.lower()
            # Only contents in section 2 is of interest
            try:
                if re.search(r"^\s*2\.", section_name):
                    if section_name.find("target") > -1:
                        slide_num = toc.cell(row, 1).text[:]
                        toc_dict["sim_targets"] = slide_num
                    elif section_name.find("mask") > -1:
                        slide_num = toc.cell(row, 1).text[:]
                        toc_dict["eye_masks"] = slide_num
                    elif section_name.find("topology") > -1:
                        slide_num = toc.cell(row, 1).text[:]
                        toc_dict["topology"] = slide_num
                # Check if end of TOC in order to end loop
                elif section_name == "":
                    break
                # Move down TOC
                row += 1 
            # In case pointer has reached end of TOC
            except IndexError:
                break

        # Convert str slide_nums to int for slide indexing
        for section, slide_nums in toc_dict.items():
            # Check if range of slide_nums
            # In case of hyphen type -
            if slide_nums.find("-") > -1:
                slide_nums = slide_nums.split("-")
                toc_dict[section] = [ int(num) - 1 for num in slide_nums ]
            # In case of hyphen type ―    
            elif slide_nums.find("\u2013") > -1:
                slide_nums = slide_nums.split("\u2013") 
                toc_dict[section] = [ int(num) - 1 for num in slide_nums ]
            # If single number
            else:
                print(slide_nums)
                toc_dict[section] = list(int(slide_nums) - 1) # Keep returned data structures consistent
        
        return toc_dict

    # def list_interfaces(self):
    #     interfaces = []
    #     toc = self.__toc.copy() # To avoid side effects
    #     start = 0
    #     end = 1

    #     # Start and end indicies matched to self.table_of_contents
    #     if len(toc["sim_targets"]) > 1:
    #         start = toc["sim_targets"][0]
    #         end += toc["sim_targets"][1]
    #     else:
    #         start = toc["sim_targets"]
    #         end += start

    #     for slide in self.pptx.slides[start:end]:
    #         try:
    #             title = slide.shapes[TITLE_SHAPE].lower()
    #             match = re.search(r"condition\s?:\s?(\w+)\b", title)
    #             # Interfaces are found in section 2 only
    #             if match:
    #                 interfaces.append(match.group(1)) # Save capture group (viz. interface)
    #         except:
    #             continue

    #     # Remove duplicates
    #     return set(interfaces)


class SimulationReport(Report):
    """
    Base class for simulation reports
    """
    __rep_types = ["si", "pi", "emc", "thermal"]

    def __init__(self, template, rep_type):
        super().__init__(template)
        self.__rep_type = rep_type

    @staticmethod
    def report_types():
        return SimulationReport.__rep_types
    
    @abstractmethod
    def __str__(self):
        return NotImplementedError

    @property
    def report_type(self):
        return self.__rep_type


class SIReport(SimulationReport):
    """
    Class for PCB signal integrity report
    """
    def __init__(self, template, interface):
        super().__init__(template=template, rep_type="SI")
        self.__interface = interface

    def __str__(self):
        return f"{self.report_type} Report for {self.interface}"

    @property
    def interface(self):
        return self.__interface
    

class PIReport(SimulationReport):
    """
    Class for PCB power integrity report
    """
    def __init__(self, template):
        super().__init__(template=template, rep_type="PI")
        self.__net_names = []

    def __str__(self):
        pass

    @property
    def net_names(self):
        return self.__net_names

    @net_names.setter
    def net_names(self, value):
        self.__net_names = value


class EMCReport(SimulationReport):
    """
    Class for PCB EMC report
    """
    def __init__(self, template):
        super().__init__(template=template, rep_type="EMC")

    def __str__(self):
        pass

class ThermalReport(SimulationReport):
    """
    Class for PCB thermal report
    """
    def __init__(self, template):
        super().__init__(template=template, rep_type="Thermal")

    def __str__(self):
        pass


# # ==============
# # Slide classes
# # ==============

# class Slide():
#     """Base class for slide in report"""
#     pass


# class CoverSlide(Slide):
#     """Class for first slide of report"""
#     pass


# class DividerSlide(Slide):
#     """Class for slide dividing sections of report"""
#     pass



# # ========================
# # SlideContent Base Class
# # ========================

# class SlideContent():
#     """Base class for images, textboxes, tables, etc. on slides"""
#     def __init__(self, wh, hasBorder, xy):
#         self.__dimensions = wh # tuple
#         self.__hasBorder = hasBorder
#         self.__position = xy # tuple
#         self.__border = {
#             "type": "",
#             "thickness": 0,
#             "style": "",
#             "color": ""
#         }
    
#     def set_border(self):
#         raise NotImplementedError


# # ==============
# # Image classes
# # ==============

# class Image(SlideContent):
#     """Base class for images on report"""
#     pass

# class TopologyImage(Image):
#     """Topology diagram"""
#     pass


# class WaveForm(Image):
#     """Waveform images for SI reports"""
#     pass


# # ==============
# # TextBox classes
# # ==============

# class TextBox(SlideContent):
#     """Base class for labels on slides"""
#     def __init__(self, ff, color, bg_color, size, hasBorder, position):
#         super().__init__(size, hasBorder, position)
#         self.__font_family = ff
#         self.__font_color = color
#         self.__bg_color = bg_color
    
    
# class Title(TextBox):
#     pass 


# class Subtitle(Title):
#     pass


# class Label(TextBox):
#     """All labels besides (sub)titles"""
#     def __init__(self, color):
#         super().__init__(color)

# class Comment(TextBox):
#    pass 


# # ==============
# # TextBox classes
# # ==============

# class Table(SlideContent):
#     """Base class for tables on report slides"""
#     def __init__(self, w, h):
#         self.__num_rows = h
#         self.__num_cols = w


# class Subtable(Table):
#     def __init__(self, w, h):
#         super().__init__(w, h)    


# class CellCollection():
#     def __init__(self, colors):
#         self.__cell_color = colors[0]
#         self.__font_color = colors[1]

# class Column(CellCollection):
#     def __init__(self, colors):
#         super().__init__(colors)


# class Row(CellCollection):
#     def __init__(self, colors):
#         super().__init__(colors)


# class Header(Row):
#     pass


if __name__ == "__main__":
    desc = """
            Weaver.py takes paths to: 
                (1) a confirmation tools report,
                (2) a Simulaton directory
                (3) a templates.txt listing templates to be used (optional)
                (4) a destination directory (optional)

            to automatically generate simulation reports 
            according to a templates listed in (3).

            For more information refer to the README.
           """
    parser = argparse.ArgumentParser(description=desc)
    # Positional args
    parser.add_argument("conf_tools", help="Path to confirmation tools for simulation reports")
    parser.add_argument("report_type", help="Type of report (SI / PI / EMC / Thermal)")
    # Optional args
    parser.add_argument("-s", "--simulation_dir", nargs=1, help="Path to simulation directory") 
    parser.add_argument("-i", "--image_dir", nargs=1, help="Path to directory of images to be included in the report(s)")

    # Retrieve args
    args = parser.parse_args()   

    # Process input from positional args    
    conf_path = args.conf_tools 
    rep_type = args.report_type.lower()
    # Verify report type
    if rep_type not in SimulationReport.report_types():
        print("ERROR: Report type is invalid")
        sys.exit(-1)
    
    # Process input from optional args
    img_dir = args.image_dir # TODO: Conditional logic for image_dir opt
    sim_dir = args.simulation_dir
    interfaces = fetch_interfaces(args.simulation_dir) if sim_dir else []

    # Make reports based on inputs and print confirmation
    weave_reports(rep_type, conf_path, interfaces)
    print(f"Weaving of report(s) for simulation type {rep_type.upper()} complete.\n")
    
    # Close program
    sleep(1)
    input("Press any key to quit.\n")